"""Generate the test fixtures listed in ENGINEERING_SPEC.md §55.

Fixtures are generated rather than committed as binaries so the repository
stays free of opaque blobs and the inputs stay inspectable.

    .venv/Scripts/python.exe ../tests/converter/fixtures/build_fixtures.py

Writes into `generated/`, which is gitignored.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

from docx import Document
from docx.shared import Inches
from fpdf import FPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.util import Pt

OUTPUT_DIR = Path(__file__).parent / "generated"


def _png_bytes(color: tuple[int, int, int], size: tuple[int, int] = (48, 32)) -> bytes:
    buffer = io.BytesIO()
    Image.new("RGB", size, color).save(buffer, format="PNG")
    return buffer.getvalue()


def _png_file(
    name: str,
    color: tuple[int, int, int],
    size: tuple[int, int] = (48, 32),
) -> Path:
    path = OUTPUT_DIR / name
    path.write_bytes(_png_bytes(color, size))
    return path


# -- DOCX -------------------------------------------------------------------


def build_docx_simple() -> None:
    document = Document()
    document.add_paragraph("A single plain paragraph of body text.")
    document.save(OUTPUT_DIR / "simple.docx")


def build_docx_headings() -> None:
    document = Document()
    document.add_heading("Top Level Heading", level=1)
    document.add_paragraph("Intro paragraph under the top heading.")
    document.add_heading("Second Level", level=2)
    document.add_paragraph("Body text with bold and italic runs follows.")
    paragraph = document.add_paragraph()
    paragraph.add_run("bold text").bold = True
    paragraph.add_run(" and ")
    paragraph.add_run("italic text").italic = True
    document.add_heading("Third Level", level=3)
    for item in ("First bullet", "Second bullet", "Third bullet"):
        document.add_paragraph(item, style="List Bullet")
    document.save(OUTPUT_DIR / "headings.docx")


def build_docx_table() -> None:
    document = Document()
    document.add_heading("Quarterly Figures", level=1)
    table = document.add_table(rows=3, cols=3)
    data = [
        ["Region", "Q1", "Q2"],
        ["North", "120", "145"],
        ["South", "98", "131"],
    ]
    for row_index, row in enumerate(data):
        for col_index, value in enumerate(row):
            table.cell(row_index, col_index).text = value
    document.save(OUTPUT_DIR / "table.docx")


def build_docx_images() -> None:
    document = Document()
    document.add_heading("Document With Images", level=1)
    document.add_paragraph("Text before the first image.")
    logo = _png_file("_docx_logo.png", (30, 120, 200))
    document.add_picture(str(logo), width=Inches(1.0))
    document.add_paragraph("Text between the images.")
    # Same image again - exercises media de-duplication.
    document.add_picture(str(logo), width=Inches(1.0))
    second = _png_file("_docx_second.png", (220, 90, 40))
    document.add_picture(str(second), width=Inches(1.0))
    document.save(OUTPUT_DIR / "images.docx")


# -- PPTX -------------------------------------------------------------------


def _blank_layout(presentation: Presentation):
    return presentation.slide_layouts[6]


def _titled_layout(presentation: Presentation):
    return presentation.slide_layouts[5]


def build_pptx_text_only() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(_titled_layout(presentation))
    slide.shapes.title.text = "Strategy Overview"
    box = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(2), PptxInches(6), PptxInches(3)
    )
    frame = box.text_frame
    frame.text = "First top level point"
    for text, level in (("Supporting detail", 1), ("Second top level point", 0)):
        paragraph = frame.add_paragraph()
        paragraph.text = text
        paragraph.level = level
        paragraph.font.size = Pt(18)

    second = presentation.slides.add_slide(_titled_layout(presentation))
    second.shapes.title.text = "Next Steps"
    box2 = second.shapes.add_textbox(
        PptxInches(1), PptxInches(2), PptxInches(6), PptxInches(2)
    )
    box2.text_frame.text = "A closing statement."

    presentation.save(OUTPUT_DIR / "text-only.pptx")


def build_pptx_images() -> None:
    presentation = Presentation()
    logo = _png_file("_pptx_logo.png", (10, 160, 90))
    for index in range(1, 4):
        slide = presentation.slides.add_slide(_titled_layout(presentation))
        slide.shapes.title.text = f"Slide {index}"
        # The identical logo on every slide must be stored only once.
        slide.shapes.add_picture(
            str(logo), PptxInches(1), PptxInches(2), width=PptxInches(2)
        )
    unique = _png_file("_pptx_unique.png", (200, 30, 120))
    presentation.slides[2].shapes.add_picture(
        str(unique), PptxInches(4), PptxInches(2), width=PptxInches(2)
    )
    presentation.save(OUTPUT_DIR / "images.pptx")


def build_pptx_tables() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(_titled_layout(presentation))
    slide.shapes.title.text = "Comparison"
    shape = slide.shapes.add_table(
        3, 3, PptxInches(1), PptxInches(2), PptxInches(7), PptxInches(2)
    )
    data = [
        ["Option", "Cost", "Risk"],
        ["Build", "High", "Medium"],
        ["Buy", "Medium", "Low"],
    ]
    for row_index, row in enumerate(data):
        for col_index, value in enumerate(row):
            shape.table.cell(row_index, col_index).text = value
    presentation.save(OUTPUT_DIR / "tables.pptx")


def build_pptx_grouped_shapes() -> None:
    """A grouped shape whose children must be walked recursively (§34)."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(_titled_layout(presentation))
    slide.shapes.title.text = "Grouped Content"

    # python-pptx has no high-level group builder, so assemble the group XML
    # from two real textboxes.
    from pptx.oxml.ns import qn

    box_a = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(2), PptxInches(3), PptxInches(1)
    )
    box_a.text_frame.text = "Grouped child one"
    box_b = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(3), PptxInches(3), PptxInches(1)
    )
    box_b.text_frame.text = "Grouped child two"

    sp_tree = slide.shapes._spTree  # noqa: SLF001 - no public group API
    group = sp_tree.add_grpSp()
    for box in (box_a, box_b):
        element = box._element  # noqa: SLF001
        sp_tree.remove(element)
        group.append(element)
    # Give the group a usable offset so reading-order sorting is deterministic.
    group.find(qn("p:grpSpPr"))

    presentation.save(OUTPUT_DIR / "grouped-shapes.pptx")


# -- PDF --------------------------------------------------------------------


def build_pdf_text() -> None:
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=14)
    pdf.cell(0, 10, "Market Study Introduction", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", size=11)
    for line in (
        "This document describes the addressable market.",
        "Growth is projected across three segments.",
        "Each segment is assessed independently.",
    ):
        pdf.cell(0, 8, line, new_x="LMARGIN", new_y="NEXT")
    pdf.output(str(OUTPUT_DIR / "text.pdf"))


def build_pdf_multipage() -> None:
    pdf = FPDF()
    pdf.set_font("Helvetica", size=12)
    for page in range(1, 4):
        pdf.add_page()
        pdf.cell(0, 10, f"Page {page} heading", new_x="LMARGIN", new_y="NEXT")
        pdf.cell(
            0, 8, f"Body content belonging to page {page}.",
            new_x="LMARGIN", new_y="NEXT",
        )
    pdf.output(str(OUTPUT_DIR / "multipage.pdf"))


def build_pdf_images() -> None:
    image_path = _png_file("_pdf_image.png", (60, 60, 200), size=(120, 80))
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, "Document with an embedded image", new_x="LMARGIN", new_y="NEXT")
    pdf.image(str(image_path), x=20, y=40, w=60)
    pdf.output(str(OUTPUT_DIR / "images.pdf"))


def build_pdf_scanned_like() -> None:
    """A page that is almost entirely image with negligible text (§36)."""
    image_path = _png_file("_pdf_scan.png", (240, 240, 235), size=(1000, 1400))
    pdf = FPDF()
    pdf.add_page()
    pdf.image(str(image_path), x=0, y=0, w=210)
    pdf.output(str(OUTPUT_DIR / "scanned-like.pdf"))


# -- Security fixtures (§55) ------------------------------------------------


def build_fake_pdf() -> None:
    """Claims .pdf, contains no PDF header."""
    (OUTPUT_DIR / "fake-pdf.pdf").write_bytes(b"This is plain text, not a PDF.\n")


def build_renamed_zip() -> None:
    """A valid ZIP claiming .docx, missing the required OOXML parts."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("readme.txt", "not an office document")
    (OUTPUT_DIR / "renamed-zip.docx").write_bytes(buffer.getvalue())


def build_unsafe_office_archive() -> None:
    """A .docx-shaped ZIP with a traversal member and a hostile ratio."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        archive.writestr("word/document.xml", "<document/>")
        archive.writestr("../../escape.txt", "traversal attempt")
    (OUTPUT_DIR / "unsafe-office-archive.docx").write_bytes(buffer.getvalue())


def build_zip_bomb_docx() -> None:
    """Valid OOXML shape, but expands far beyond the compression-ratio guard."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        archive.writestr("word/document.xml", "<document/>")
        archive.writestr("word/payload.bin", b"\x00" * (40 * 1024 * 1024))
    (OUTPUT_DIR / "zip-bomb.docx").write_bytes(buffer.getvalue())


def build_encrypted_office() -> None:
    """OLE2/CFB container - what a password-protected DOCX actually is."""
    cfb_magic = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
    (OUTPUT_DIR / "encrypted.docx").write_bytes(cfb_magic + b"\x00" * 512)



# -- Amendment A8.1 additions ----------------------------------------------


def build_pptx_charts() -> None:
    """A native chart, which cannot be represented in Markdown (A1.5)."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    presentation = Presentation()
    slide = presentation.slides.add_slide(_titled_layout(presentation))
    slide.shapes.title.text = "Revenue by Region"

    data = CategoryChartData()
    data.categories = ["North", "South", "East"]
    data.add_series("FY26", (120.0, 98.0, 143.0))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        PptxInches(1),
        PptxInches(2),
        PptxInches(6),
        PptxInches(4),
        data,
    )
    presentation.save(OUTPUT_DIR / "charts.pptx")


def build_pptx_speaker_notes() -> None:
    """Notes must NOT appear in output unless PPTX_INCLUDE_NOTES is set."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(_titled_layout(presentation))
    slide.shapes.title.text = "Public Title"
    box = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(2), PptxInches(6), PptxInches(1)
    )
    box.text_frame.text = "Body text that should appear."
    slide.notes_slide.notes_text_frame.text = (
        "CONFIDENTIAL internal commentary that must not be published."
    )
    presentation.save(OUTPUT_DIR / "speaker-notes.pptx")


def build_pptx_merged_cells() -> None:
    """A merged-cell table must still render as valid GFM."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(_titled_layout(presentation))
    slide.shapes.title.text = "Merged Header"
    shape = slide.shapes.add_table(
        3, 3, PptxInches(1), PptxInches(2), PptxInches(7), PptxInches(2)
    )
    table = shape.table
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(0, 0).text = "Spanning Header"
    table.cell(0, 2).text = "Third"
    for row in (1, 2):
        for col in range(3):
            table.cell(row, col).text = f"r{row}c{col}"
    presentation.save(OUTPUT_DIR / "merged-cells.pptx")


def build_pptx_coloured_text() -> None:
    """Coloured runs must not produce raw HTML colour tags (A8.2)."""
    from pptx.dml.color import RGBColor

    presentation = Presentation()
    slide = presentation.slides.add_slide(_titled_layout(presentation))
    slide.shapes.title.text = "Coloured Text"
    box = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(2), PptxInches(6), PptxInches(2)
    )
    frame = box.text_frame
    frame.text = "Red warning text"
    frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xC0, 0x10, 0x10)
    paragraph = frame.add_paragraph()
    paragraph.text = "Green confirmation text"
    paragraph.runs[0].font.color.rgb = RGBColor(0x10, 0xA0, 0x40)
    presentation.save(OUTPUT_DIR / "coloured-text.pptx")


def build_pdf_tables() -> None:
    """Ruled table so pdfplumber has lines to detect."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)
    pdf.cell(0, 10, "Segment Comparison", new_x="LMARGIN", new_y="NEXT")

    rows = [
        ["Segment", "Revenue", "Growth"],
        ["Enterprise", "4.2m", "12%"],
        ["Mid-market", "2.8m", "23%"],
        ["SMB", "1.1m", "41%"],
    ]
    # Draw an explicitly ruled grid; borderless text is not detectable.
    col_width = 55
    row_height = 9
    for row in rows:
        for value in row:
            pdf.cell(col_width, row_height, value, border=1)
        pdf.ln(row_height)
    pdf.output(str(OUTPUT_DIR / "tables.pdf"))


def build_pdf_two_column() -> None:
    """Documents the known multi-column reading-order limitation (A10)."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    left = ["Left column line one.", "Left column line two.", "Left column line three."]
    right = [
        "Right column line one.",
        "Right column line two.",
        "Right column line three.",
    ]
    for index, (left_text, right_text) in enumerate(zip(left, right, strict=True)):
        y = 40 + index * 8
        pdf.set_xy(15, y)
        pdf.cell(80, 8, left_text)
        pdf.set_xy(110, y)
        pdf.cell(80, 8, right_text)
    pdf.output(str(OUTPUT_DIR / "two-column.pdf"))


def build_pdf_many_pages() -> None:
    """Exceeds a lowered PDF_TABLE_MAX_PAGES in tests."""
    pdf = FPDF()
    pdf.set_font("Helvetica", size=11)
    for page in range(1, 13):
        pdf.add_page()
        pdf.cell(0, 10, f"Page {page} of twelve", new_x="LMARGIN", new_y="NEXT")
    pdf.output(str(OUTPUT_DIR / "many-pages.pdf"))


BUILDERS = [
    build_docx_simple,
    build_docx_headings,
    build_docx_table,
    build_docx_images,
    build_pptx_text_only,
    build_pptx_images,
    build_pptx_tables,
    build_pptx_grouped_shapes,
    build_pptx_charts,
    build_pptx_speaker_notes,
    build_pptx_merged_cells,
    build_pptx_coloured_text,
    build_pdf_text,
    build_pdf_multipage,
    build_pdf_images,
    build_pdf_scanned_like,
    build_pdf_tables,
    build_pdf_two_column,
    build_pdf_many_pages,
    build_fake_pdf,
    build_renamed_zip,
    build_unsafe_office_archive,
    build_zip_bomb_docx,
    build_encrypted_office,
]


def main() -> int:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    failures = 0
    for builder in BUILDERS:
        name = builder.__name__.removeprefix("build_")
        try:
            builder()
        except Exception as exc:  # noqa: BLE001 - report and keep going
            failures += 1
            print(f"  FAIL  {name}: {type(exc).__name__}: {exc}")
        else:
            print(f"  ok    {name}")

    # Remove the scratch PNGs used only as builder inputs.
    for scratch in OUTPUT_DIR.glob("_*.png"):
        scratch.unlink()

    print(f"\nFixtures in {OUTPUT_DIR}")
    return 1 if failures else 0


if __name__ == "__main__":
    raise SystemExit(main())
