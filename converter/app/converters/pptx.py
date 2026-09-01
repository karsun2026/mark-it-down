"""PPTX -> Markdown via python-pptx (ENGINEERING_SPEC.md §34).

No PowerPoint dependency, no AI. Shapes are emitted in approximate reading
order (top, then left), grouped shapes are walked recursively, native tables
become GitHub-style Markdown tables, and anything unsupported produces a
warning rather than failing the job.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.config import settings
from app.converters.base import BaseConverter, ConversionResult, sniff_image_extension
from app.errors import ConversionError, ErrorCode

logger = logging.getLogger(__name__)

# Amendment A1.5 warning taxonomy. One warning per distinct unsupported class
# per slide, worded for a non-technical reader and naming the slide so the user
# knows where to look. Keys are MSO_SHAPE_TYPE members.
#
# §39/§47 apply: warning text never contains document text, shape text or paths.
_UNSUPPORTED_WARNINGS: dict[Any, str] = {
    MSO_SHAPE_TYPE.CHART: "Slide {n} contains a chart. Chart data is not converted.",
    MSO_SHAPE_TYPE.DIAGRAM: (
        "Slide {n} contains SmartArt. SmartArt is not converted."
    ),
    MSO_SHAPE_TYPE.MEDIA: (
        "Slide {n} contains embedded media. Media is not converted."
    ),
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: (
        "Slide {n} contains an embedded object. It is not converted."
    ),
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: (
        "Slide {n} contains a linked object. It is not converted."
    ),
    MSO_SHAPE_TYPE.LINE: (
        "Slide {n} contains connectors or diagram lines. Spatial relationships "
        "may be lost."
    ),
}

# Vector image formats we deliberately leave untouched. Converting WMF needs an
# extra native dependency (A7.1 prohibits `wand`), so these warn instead.
_UNCONVERTED_IMAGE_EXTENSIONS = frozenset({"wmf", "emf"})


# Alt text that is really just the source image's filename, which tools set by
# default and which must not reach the output.
_FILENAME_LIKE = re.compile(
    r"^[\w \-.()]+\.(png|jpe?g|gif|bmp|tiff?|webp|emf|wmf|svg)$",
    re.IGNORECASE,
)


def _escape_cell(text: str) -> str:
    """Make text safe inside a Markdown table cell."""
    return text.replace("|", "\\|").replace("\n", " ").strip()


def _shape_sort_key(shape: Any) -> tuple[int, int]:
    """Approximate reading order: top coordinate, then left (§34).

    Shapes with no position (rare, but possible for placeholders) sort last.
    """
    top = shape.top if shape.top is not None else 10**9
    left = shape.left if shape.left is not None else 10**9
    return (top, left)


class PptxConverter(BaseConverter):
    source_label = "pptx"

    def convert(self, source_path: Path) -> ConversionResult:
        try:
            presentation = Presentation(str(source_path))
        except Exception as exc:  # noqa: BLE001 - library raises many types
            raise ConversionError(
                ErrorCode.CONVERSION_FAILED,
                internal_detail=f"python-pptx open failed: {type(exc).__name__}",
            ) from exc

        blocks: list[str] = []
        slide_count = 0

        for index, slide in enumerate(presentation.slides, start=1):
            slide_count = index
            blocks.extend(self._convert_slide(slide, index))

        markdown_path = self._write_markdown(blocks)

        return ConversionResult(
            markdown_path=markdown_path,
            media_dir=self.workspace.media_dir,
            pages_or_slides=slide_count,
            media_count=self.media.count,
            warnings=self.warnings,
        )

    # -- slide handling ----------------------------------------------------

    def _convert_slide(self, slide: Any, index: int) -> list[str]:
        title = self._slide_title(slide)
        heading = f"## Slide {index} — {title}" if title else f"## Slide {index}"

        blocks: list[str] = []
        if index > 1:
            blocks.append("---")
        blocks.append(heading)

        # python-pptx returns a fresh proxy object on each `.title` access, so
        # identity comparison never matches. Compare the stable shape id.
        title_id = self._title_shape_id(slide)
        for shape in sorted(slide.shapes, key=_shape_sort_key):
            # The title is already the heading; don't repeat it in the body.
            if title_id is not None and getattr(shape, "shape_id", None) == title_id:
                continue
            blocks.extend(self._convert_shape(shape, index))

        # A1.3: speaker notes are excluded by default. They routinely carry
        # internal commentary the uploader does not expect to appear in a file
        # they then share, so inclusion is opt-in rather than the default.
        if settings.pptx_include_notes:
            notes = self._slide_notes(slide)
            if notes:
                blocks.append(f"**Speaker notes:** {notes}")

        return blocks

    def _title_shape_id(self, slide: Any) -> int | None:
        try:
            title_shape = slide.shapes.title
        except (AttributeError, ValueError):
            return None
        if title_shape is None:
            return None
        return getattr(title_shape, "shape_id", None)

    def _slide_title(self, slide: Any) -> str:
        try:
            title_shape = slide.shapes.title
        except (AttributeError, ValueError):
            return ""
        if title_shape is None or not title_shape.has_text_frame:
            return ""
        return " ".join(title_shape.text_frame.text.split()).strip()

    def _slide_notes(self, slide: Any) -> str:
        try:
            if not slide.has_notes_slide:
                return ""
            frame = slide.notes_slide.notes_text_frame
        except (AttributeError, ValueError):
            return ""
        if frame is None:
            return ""
        return " ".join(frame.text.split()).strip()

    # -- shape handling ----------------------------------------------------

    def _convert_shape(self, shape: Any, slide_index: int) -> list[str]:
        shape_type = getattr(shape, "shape_type", None)

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            blocks: list[str] = []
            try:
                children = sorted(shape.shapes, key=_shape_sort_key)
            except (AttributeError, ValueError):
                return []
            for child in children:
                blocks.extend(self._convert_shape(child, slide_index))
            return blocks

        if getattr(shape, "has_table", False):
            return self._convert_table(shape)

        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self._convert_picture(shape, slide_index)

        if shape_type in _UNSUPPORTED_WARNINGS:
            # `warn` de-duplicates, giving A1.5's one-per-class-per-slide rule.
            self.warn(_UNSUPPORTED_WARNINGS[shape_type].format(n=slide_index))
            return []

        if getattr(shape, "has_text_frame", False):
            return self._convert_text_frame(shape)

        return []

    def _convert_text_frame(self, shape: Any) -> list[str]:
        lines: list[str] = []
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if not text:
                continue
            # PowerPoint indent levels map onto nested Markdown bullets.
            level = paragraph.level or 0
            if level > 0:
                lines.append(f"{'  ' * level}- {text}")
            else:
                lines.append(f"- {text}")

        if not lines:
            return []

        # A single unindented line reads better as a paragraph than a bullet.
        if len(lines) == 1 and lines[0].startswith("- "):
            return [lines[0][2:]]
        return ["\n".join(lines)]

    def _convert_table(self, shape: Any) -> list[str]:
        try:
            table = shape.table
            rows = list(table.rows)
        except (AttributeError, ValueError):
            return []
        if not rows:
            return []

        matrix: list[list[str]] = []
        for row in rows:
            matrix.append([_escape_cell(cell.text) for cell in row.cells])

        width = max(len(row) for row in matrix)
        matrix = [row + [""] * (width - len(row)) for row in matrix]

        header, *body = matrix
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * width) + " |",
        ]
        lines.extend("| " + " | ".join(row) + " |" for row in body)
        return ["\n".join(lines)]

    def _convert_picture(self, shape: Any, slide_index: int) -> list[str]:
        try:
            image = shape.image
            data = image.blob
        except Exception as exc:  # noqa: BLE001 - unsupported/corrupt media
            self.warn(
                f"Slide {slide_index} contains an image that could not be "
                "extracted."
            )
            logger.info("pptx image extraction failed: %s", type(exc).__name__)
            return []

        # A1.5: vector metafiles are left unconverted by design rather than
        # pulling in a native conversion dependency.
        image_ext = str(getattr(image, "ext", "") or "").lower().lstrip(".")
        if image_ext in _UNCONVERTED_IMAGE_EXTENSIONS:
            self.warn(
                f"Slide {slide_index} contains a {image_ext.upper()} image. "
                "It was left unconverted."
            )
            return []

        extension = sniff_image_extension(
            data, fallback=f".{getattr(image, 'ext', 'png')}"
        )
        index = self.media.count + 1
        name = f"slide-{slide_index:03d}-image-{index:03d}{extension}"
        relative = self.media.write(data, name)

        alt = self._picture_alt_text(shape) or f"Slide {slide_index} image"
        return [f"![{alt}]({relative})"]

    def _picture_alt_text(self, shape: Any) -> str:
        """Authored alt text, if there is any worth using.

        PowerPoint stores alt text in the shape's `descr` attribute, but tools
        (python-pptx included) default it to the source image's filename. That
        is not descriptive and leaks the author's local filename into the
        output, so filename-shaped values are discarded.
        """
        try:
            descr = shape._element._nvXxPr.cNvPr.get("descr")  # noqa: SLF001
        except (AttributeError, KeyError, ValueError):
            return ""

        text = " ".join((descr or "").split()).strip()
        if not text:
            return ""
        if _FILENAME_LIKE.match(text):
            return ""
        return text
