"""Generate near-limit source documents for the §57 release test.

    .venv/Scripts/python.exe scripts/make_large_fixtures.py --mb 95 --format pptx

§57 requires a real file in the 95-100 MB band. The content matters as much as
the size:

  * Every embedded image is **unique random noise**. Identical images would be
    collapsed by the content-addressed de-duplication in `MediaWriter`, so a
    deck built from one repeated image would produce a tiny output tree and
    prove nothing about behaviour at the ceiling.
  * Noise is incompressible, so the archive's compression ratio stays near 1
    and the file does not trip the §30 ZIP-bomb guard, which is testing a
    different failure mode.

Output goes to `tests/converter/fixtures/large/`, which is gitignored — these
files are far too big to commit.
"""

from __future__ import annotations

import argparse
import io
import os
import sys
import time
from pathlib import Path

REPO = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(REPO / "converter"))

OUTPUT_DIR = REPO / "tests" / "converter" / "fixtures" / "large"

MEGABYTE = 1024 * 1024
# 600x600 RGB noise is ~1.03 MB raw and PNG-encodes to about the same, since
# random data does not compress.
NOISE_EDGE = 600
BYTES_PER_NOISE_IMAGE = NOISE_EDGE * NOISE_EDGE * 3


def _noise_png() -> bytes:
    """One unique, incompressible PNG."""
    from PIL import Image

    image = Image.frombytes(
        "RGB", (NOISE_EDGE, NOISE_EDGE), os.urandom(BYTES_PER_NOISE_IMAGE)
    )
    buffer = io.BytesIO()
    # compress_level=0 keeps generation fast and the bytes incompressible.
    image.save(buffer, format="PNG", compress_level=0)
    return buffer.getvalue()


def _progress(done: int, total: int, started: float) -> None:
    elapsed = time.perf_counter() - started
    print(
        f"\r  {done}/{total} images  {elapsed:6.1f}s",
        end="",
        flush=True,
    )


def build_pptx(target_bytes: int) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    layout = presentation.slide_layouts[5]
    count = max(1, target_bytes // BYTES_PER_NOISE_IMAGE)
    started = time.perf_counter()

    for index in range(1, count + 1):
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = f"Section {index}"
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.4), Inches(4), Inches(0.6)
        )
        box.text_frame.text = f"Narrative body text for section {index}."
        slide.shapes.add_picture(
            io.BytesIO(_noise_png()), Inches(0.5), Inches(2.2), width=Inches(4)
        )
        if index % 10 == 0 or index == count:
            _progress(index, count, started)

    print()
    path = OUTPUT_DIR / "near-limit.pptx"
    presentation.save(str(path))
    return path


def build_pdf(target_bytes: int) -> Path:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_font("Helvetica", size=11)
    count = max(1, target_bytes // BYTES_PER_NOISE_IMAGE)
    started = time.perf_counter()

    for index in range(1, count + 1):
        pdf.add_page()
        pdf.cell(0, 10, f"Page {index} heading", new_x="LMARGIN", new_y="NEXT")
        pdf.cell(
            0, 8, f"Body content for page {index}.", new_x="LMARGIN", new_y="NEXT"
        )
        # fpdf2 accepts a file-like image.
        pdf.image(io.BytesIO(_noise_png()), x=20, y=40, w=120)
        if index % 10 == 0 or index == count:
            _progress(index, count, started)

    print()
    path = OUTPUT_DIR / "near-limit.pdf"
    pdf.output(str(path))
    return path


def build_docx(target_bytes: int) -> Path:
    from docx import Document
    from docx.shared import Inches

    document = Document()
    count = max(1, target_bytes // BYTES_PER_NOISE_IMAGE)
    started = time.perf_counter()

    for index in range(1, count + 1):
        document.add_heading(f"Section {index}", level=2)
        document.add_paragraph(f"Narrative body text for section {index}.")
        document.add_picture(io.BytesIO(_noise_png()), width=Inches(4))
        if index % 10 == 0 or index == count:
            _progress(index, count, started)

    print()
    path = OUTPUT_DIR / "near-limit.docx"
    document.save(str(path))
    return path


BUILDERS = {"pptx": build_pptx, "pdf": build_pdf, "docx": build_docx}


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--mb", type=int, default=95, help="target source size in MB (default 95)"
    )
    parser.add_argument(
        "--format",
        choices=sorted(BUILDERS) + ["all"],
        default="all",
        help="which fixture to build",
    )
    args = parser.parse_args()

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    formats = sorted(BUILDERS) if args.format == "all" else [args.format]
    target = args.mb * MEGABYTE

    for name in formats:
        print(f"building {name} (~{args.mb} MB)")
        started = time.perf_counter()
        path = BUILDERS[name](target)
        size = path.stat().st_size
        elapsed = time.perf_counter() - started
        band = "OK" if 95 * MEGABYTE <= size <= 100 * MEGABYTE else "OUT OF BAND"
        print(
            f"  {path.name}: {size / MEGABYTE:.1f} MB  "
            f"({elapsed:.0f}s)  [{band} for §57's 95-100 MB requirement]"
        )

    print(f"\nFixtures in {OUTPUT_DIR}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
